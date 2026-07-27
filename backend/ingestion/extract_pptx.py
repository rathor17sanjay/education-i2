"""PowerPoint (.pptx) text extraction. Same extract(path) -> str contract as
extract_pdf.py -- "--- slide N ---" markers mirror extract_pdf.py's
"--- page N ---" convention, so retrieved-chunk provenance reads the same
way regardless of source type.
"""

from pptx import Presentation


def _shape_text(shape) -> str:
    if shape.has_text_frame:
        return shape.text_frame.text
    if shape.has_table:
        rows = [
            "\t".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows
        ]
        return "\n".join(rows)
    return ""


def extract(path: str) -> str:
    prs = Presentation(path)

    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts = [_shape_text(shape) for shape in slide.shapes]
        texts = [t for t in texts if t.strip()]
        if texts:
            slides_text.append(f"--- slide {i + 1} ---\n" + "\n".join(texts))

    return "\n\n".join(slides_text)
